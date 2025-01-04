import boto3
from django.conf import settings
from botocore.exceptions import ClientError, NoCredentialsError
from pptx import Presentation
from io import BytesIO
import base64
import json

class S3ServiceError(Exception):
    """Custom exception for S3Service errors"""
    pass

class S3Service:
    _instance = None
    _s3_client = None

    @classmethod
    def get_instance(cls):
        if cls._instance is None:
            cls._instance = cls()
        return cls._instance

    def __init__(self):
        if self._s3_client is None:
            try:
                # Debug: Print to verify credentials are loaded
                # print(f"Region: {settings.AWS_S3_REGION_NAME}")
                # print(f"Access Key: {settings.AWS_ACCESS_KEY_ID[:4]}...")  # Print first 4 chars only
                # print(f"Secret Access Key: {settings.AWS_SECRET_ACCESS_KEY[:4]}...")  # Print first 4 chars only
                # print(f"Bucket Name: {settings.AWS_STORAGE_BUCKET_NAME}")
                
                session = boto3.Session(
                    aws_access_key_id=settings.AWS_ACCESS_KEY_ID,
                    aws_secret_access_key=settings.AWS_SECRET_ACCESS_KEY,
                    region_name=settings.AWS_S3_REGION_NAME
                )
                self._s3_client = session.client('s3')
            except NoCredentialsError:
                raise S3ServiceError("AWS credentials not found")
            except ClientError as e:
                raise S3ServiceError(f"AWS client error: {str(e)}")
            except Exception as e:
                raise S3ServiceError(f"Failed to initialize S3 client: {str(e)}")

    def get_presigned_url(self, slide_id):
        """
        Generates a pre-signed URL for accessing a PNG file in S3.
        
        This function does the following:
        1. Lists objects in S3 bucket that match the given png_id prefix
        2. Gets the full object key of the first matching file
        3. Generates a pre-signed URL that allows temporary access to the file
        
        Args:
            png_id (str): The ID/prefix of the PNG file to generate URL for
            
        Returns:
            str: Pre-signed URL valid for 1 hour if successful
            None: If there's an error or file not found
            
        The pre-signed URL allows temporary access to private S3 objects
        without requiring AWS credentials. It expires after 1 hour (3600 seconds).
        """
        try:
            response = self._s3_client.list_objects_v2(
                Bucket=settings.AWS_STORAGE_BUCKET_NAME,
                Prefix=f"slide_images/images/{slide_id}"
            )
            
            if 'Contents' not in response or response['KeyCount'] == 0:
                print(f"No objects found with prefix: {slide_id}")
                return None
            
            # Get the first matching object's key
            object_key = response['Contents'][0]['Key']
            # object_key = "slide_images/images/070f1910-a717-4600-98fc-a569a5433cb7.png"
            
            presigned_url = self._s3_client.generate_presigned_url('get_object',
                Params={
                    'Bucket': settings.AWS_STORAGE_BUCKET_NAME,
                    'Key': object_key
                },
                ExpiresIn=3600
            )
            return presigned_url
            
        except ClientError as e:
            print(f"Error: {e}")
            return None
    
    def extract_slides_data_from_s3(self, file_name):
        try:
            ids = ["070f1910-a717-4600-98fc-a569a5433cb7",
                    # "f343d895-fc63-4d9f-8610-6ab9da6662fb",
                    # "1a7628d9-694c-4202-b942-c265887d8ac1",
                    "19c9f5d5-d686-4628-8b60-4b170c839f14",
                    # "d06a5c9a-1937-4e84-9051-c54fb4b9c1c7"
                    ]
            try:
                response = self._s3_client.get_object(
                    Bucket=settings.AWS_STORAGE_BUCKET_NAME,
                    Key='slide_metadata/metadata_json/metadata.json'
                )
            except ClientError as e:
                print(f"Error retrieving metadata.json: {e}")
                return None
            
            json_data = json.loads(response['Body'].read().decode('utf-8'))
            json_ids = json_data.get('slide_ids', {})
            # Get matching IDs between hardcoded list and JSON data
            matching_idxs = [i for i, id in enumerate(json_ids.values()) if id in ids]
            print(f"Matching IDs: {matching_idxs}")

            response = self._s3_client.get_object(
                Bucket=settings.AWS_STORAGE_BUCKET_NAME,
                Key=file_name
            )
            file_stream = response['Body']
            pptx_file = BytesIO(file_stream.read())
            presentation = Presentation(pptx_file)

            slides = []
            for i, slide in enumerate(presentation.slides):
                if i in matching_idxs:
                    slide_content = {
                        "texts": [shape.text for shape in slide.shapes if shape.has_text_frame],
                        "images": [base64.b64encode(shape.image.blob).decode('utf-8') 
                                for shape in slide.shapes if hasattr(shape, "image")]
                    }
                    slides.append(slide_content)
            print(f'Slides End: {len(slides)}')
            return slides
        except Exception as e:
            print(f"Error: {e}")
            return None

    def upload_file(self, file_obj, key):
        try:
            return self._s3_client.upload_fileobj(
                file_obj,
                settings.AWS_STORAGE_BUCKET_NAME,
                key
            )
        except ClientError as e:
            print(f"Error: {e}")
            return None

    def convert_pptx_to_base64(self, slides):
        try:
            ids = [slide['slide_id'] for slide in slides if slide['file_id'] == 'mckinsey.pptx']
            try:
                response = self._s3_client.get_object(
                    Bucket=settings.AWS_STORAGE_BUCKET_NAME,
                    Key='slide_metadata/metadata_json/metadata.json'
                )
            except ClientError as e:
                print(f"Error retrieving metadata.json: {e}")
                return None
            
            json_data = json.loads(response['Body'].read().decode('utf-8'))
            json_ids = json_data.get('slide_ids', {})
            # Get matching IDs between hardcoded list and JSON data
            matching_idxs = [i for i, id in enumerate(json_ids.values()) if id in ids]
            print(f"Matching IDXs: {matching_idxs}")

            # Get the file from S3
            response = self._s3_client.get_object(
                Bucket=settings.AWS_STORAGE_BUCKET_NAME,
                Key='mckinsey.pptx'
            )
            
            # Read the file into memory
            file_stream = response['Body']
            file_data = file_stream.read()
            
            # Convert to base64 and create data URL
            base64_data = base64.b64encode(file_data).decode('utf-8')
            
            # Extract just the base64 data after the data URL prefix
            data_url = f"data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{base64_data}"
            start_idx = data_url.index("base64,") + 7
            base64_only = data_url[start_idx:]

            pptx_file = BytesIO(file_data)
            presentation = Presentation(pptx_file)

            slide_ids = []
            for i, slide in enumerate(presentation.slides):
                if i in matching_idxs:
                    slide_ids.append(str(slide.slide_id))  # Get the slide ID in string format
            print(f'Slides End: {slide_ids}')
            
            # Return list with single base64 string since we're getting the whole file
            return base64_only, slide_ids
            
        except Exception as e:
            print(f"Error converting slides to base64: {e}")
            return None
    
    def slides_indexex_from_metadata(self, file, slide_ids):
        try:
            print(f'File: {file}, Slide ids: {slide_ids}')
            response = self._s3_client.get_object(
                Bucket=settings.AWS_STORAGE_BUCKET_NAME,
                Key='slide_metadata/metadata_json/metadata.json'
            )
            json_data = json.loads(response['Body'].read().decode('utf-8'))
            json_ids = json_data.get('slide_ids', {})
            matching_idxs = [i for i, id in enumerate(json_ids.values()) if id in slide_ids]
            return file, matching_idxs
        except Exception as e:
            print(f"Error retrieving metadata.json: {e}")
            return None

    def convert_selected_slides_to_base64(self, selected_data):
        """
        Converts selected slides from a PowerPoint file to base64 string.
        
        Args:
            selected_data (list): List of slide data to include
            
        Returns:
            str: Base64 encoded string of the PowerPoint with selected slides
        """
        try:
            print(f'Selected data: {selected_data}')
            # Group slides by file_id
            slides_by_file = {}
            for slide in selected_data:
                file_id = slide['file_id']
                slide_id = slide['slide_id']
                if file_id not in slides_by_file:
                    slides_by_file[file_id] = []
                slides_by_file[file_id].append(slide_id)

            print(f'Slides by file: {slides_by_file}')

            # Create new presentation
            target_prs = Presentation()

            for file_id, slide_ids in slides_by_file.items():
                file, matching_idxs = self.slides_indexex_from_metadata(file_id, slide_ids)
                print(f'File: {file}, Matching idxs: {matching_idxs}')
                # Get the file from S3
                response = self._s3_client.get_object(
                    Bucket=settings.AWS_STORAGE_BUCKET_NAME,
                    Key=file
                )
            
                # Read the file into memory
                file_stream = response['Body']
                pptx_file = BytesIO(file_stream.read())
                source_prs = Presentation(pptx_file)
            
                # Copy selected slides
                for i, slide in enumerate(source_prs.slides):
                    if i in matching_idxs:
                        print(f'Slide: {i}, Matching idx: {matching_idxs}')
                        # Get the layout from source
                        source_layout = slide.slide_layout
                        # Find matching layout in target
                        target_layout = None
                        for layout in target_prs.slide_layouts:
                            if layout.name == source_layout.name:
                                target_layout = layout
                                break
                        if not target_layout:
                            # If no matching layout found, use first available
                            target_layout = target_prs.slide_layouts[0]
                        
                        # Add new slide with matching layout
                        target_slide = target_prs.slides.add_slide(target_layout)
                        
                        # Copy all shapes
                        for shape in slide.shapes:
                            elem = shape.element
                            target_slide.shapes._spTree.insert_element_before(elem)
            
            # Save the presentation to bytes
            output = BytesIO()
            target_prs.save(output)
            output.seek(0)
            
            # Convert to base64
            base64_string = base64.b64encode(output.getvalue()).decode('utf-8')
            
            return base64_string
            
        except Exception as e:
            print(f"Error converting slides to base64: {e}")
            return None
